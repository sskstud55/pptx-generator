from flask import Flask, render_template, request, jsonify, send_from_directory
from pptx import Presentation
import os

app = Flask(__name__)

# Directory to store generated presentations
PRESENTATIONS_DIR = 'presentations'
if not os.path.exists(PRESENTATIONS_DIR):
    os.makedirs(PRESENTATIONS_DIR)

def create_presentation(title):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title_placeholder = slide.shapes.title
    title_placeholder.text = title

    filename = f"{title.replace(' ', '_').replace('/', '_')}.pptx"
    filepath = os.path.join(PRESENTATIONS_DIR, filename)
    prs.save(filepath)
    return filename

@app.route("/")
def index():
    return render_template("index.html")

@app.route("/chat", methods=["POST"])
def chat():
    message = request.form.get("message")
    if "create presentation" in message.lower():
        title = message.lower().replace("create presentation", "").strip()
        if title:
            filename = create_presentation(title)
            response = f"Presentation '{title}' created. <a href='/download/{filename}'>Download</a>"
        else:
            response = "Please provide a title for the presentation."
    else:
        response = f"You said: {message}"
    return jsonify({"response": response})

@app.route("/download/<filename>")
def download(filename):
    return send_from_directory(PRESENTATIONS_DIR, filename, as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True)
